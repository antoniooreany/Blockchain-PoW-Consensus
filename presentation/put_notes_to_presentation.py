#   Copyright (c) 2024, Anton Gorshkov
#   All rights reserved.
#   This code is for a main.py and its unit tests.
#   For any questions or concerns, please contact Anton Gorshkov at antoniooreany@gmail.com
import os

from pptx import Presentation

from presentation.create_enhanced_presentation import output_path

# Load the uploaded presentation
presentation_path = "./THESIS_ser - Copy.pptx"
presentation = Presentation(presentation_path)

# Define the notes for the slides
notes_content = [
    "Introduce yourself and the topic. Mention your supervisors, Prof. Dr. Martin Severin and Prof. Dr. Philipp Graf. State the objective of your thesis: understanding Proof-of-Work and optimizing blockchain performance.",
    "Provide a brief roadmap of your presentation. Highlight key sections: Blockchain fundamentals, methodology, implementation, results, and conclusion. Mention that you'll discuss the scope and limitations at the end.",
    "Explain the core research questions: How does dynamic difficulty adjustment influence block mining time? How does floating-point bit difficulty improve performance and stability? Emphasize their importance in enhancing blockchain systems.",
    "Define blockchain and key terms (e.g., block, nonce, hash, miner). Describe the significance of Proof-of-Work in ensuring security and decentralization. Briefly touch on the challenges like energy consumption and scalability.",
    "Outline the framework of your research: Floating-point bit difficulty for better granularity. Dynamic difficulty adjustment for stability. Mention the tools you used (Python libraries like hashlib, numpy, and matplotlib). State that this methodology underpins the prototype you developed.",
    "Walk through the technical details: Block structure: index, data, timestamp, nonce, previous hash. Blockchain functionality: validation, dynamic difficulty adjustment, and mining process. Explain how you simulated the system using Python. Optionally, show UML diagrams for better visualization.",
    "Summarize your findings: Floating-point bit difficulty stabilizes block creation times. Dynamic difficulty adjustment prevents abrupt shifts in mining difficulty. Present supporting data with graphs or charts showing trends in mining time and difficulty.",
    "Clarify the limitations of your study: Focused on a single-miner environment. Simplifications in transaction diversity and miner incentives. Highlight the scope for future work: Extending to multi-miner setups. Exploring hybrid consensus mechanisms.",
    "Recap your key contributions: Developed a stable blockchain prototype. Validated the advantages of fractional difficulty and dynamic adjustment. Reiterate the potential applications of your findings in real-world scenarios.",
    "Thank your supervisors for their guidance. Acknowledge colleagues and family for their support. Show appreciation for the resources provided by the university.",
    "Invite questions from the audience. Prepare to discuss both theoretical and practical aspects of your research. Emphasize your readiness to clarify what the reviewers find unclear."
]

# Add notes to each slide
for i, slide in enumerate(presentation.slides):
    if i < len(notes_content):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_content[i]


# Define the absolute path to save the presentation
output_presentation_path = os.path.abspath("./THESIS_with_Notes.pptx")

# Save the modified presentation
presentation.save(output_presentation_path)

print(f"Presentation saved at: {output_presentation_path}")
