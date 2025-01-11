#   Copyright (c) 2024, Anton Gorshkov
#   All rights reserved.
#   This code is for a main.py and its unit tests.
#   For any questions or concerns, please contact Anton Gorshkov at antoniooreany@gmail.com

# from presentation.create_presentation import create_presentation

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
presentation = Presentation()

# Define slide content
slides_content = [
    {"title": "Blockchain Technology", "content": [
        "Bachelor Thesis: Consensus Mechanism Proof-of-Work",
        "Author: Anton Gorshkov",
        "Supervisors: Prof. Dr. rer. nat. Martin Severin, Prof. Dr.-Ing. Philipp Graf",
        "Submission Date: 30.11.2024"]},
    {"title": "Introduction", "content": [
        "Motivation: Blockchain as a disruptive innovation.",
        "Problem Statement: Challenges in energy, scalability, and centralization.",
        "Goal: Develop a blockchain prototype addressing these challenges."]},
    {"title": "Research Questions", "content": [
        "How does floating-point bit difficulty improve performance and stability?",
        "What are the trade-offs between simplicity and scalability?",
        "How does dynamic difficulty adjustment influence block mining time?"]},
    {"title": "Blockchain Fundamentals", "content": [
        "Key Terms: Blockchain, Proof-of-Work, Miner, Nonce, Hash.",
        "Blockchain ensures security through cryptographic hashes and consensus."]},
    {"title": "Methods", "content": [
        "Prototype Implementation:",
        "- Floating-point bit difficulty for granular control.",
        "- Dynamic difficulty adjustment for stability.",
        "Tools: Python libraries like hashlib, numpy, and matplotlib."]},
    {"title": "Implementation", "content": [
        "Core Components: Block and Blockchain classes.",
        "Features: Dynamic difficulty adjustment, Proof-of-Work integration.",
        "Includes UML diagrams and example code snippets."]},
    {"title": "Results", "content": [
        "Findings:",
        "- Stable block creation times with fractional difficulty.",
        "- Dynamic adjustments to prevent abrupt shifts.",
        "Visuals: Graphs on mining times and difficulty trends."]},
    {"title": "Limitations", "content": [
        "Single-miner setup limits scalability insights.",
        "Simplifications in transactions and miner environment."]},
    {"title": "Conclusion", "content": [
        "Achievements: Developed and tested a blockchain prototype.",
        "Insights: Demonstrated benefits of fractional difficulty adjustments.",
        "Future Work: Multi-miner environments, hybrid consensus mechanisms."]},
    {"title": "Acknowledgments", "content": [
        "Thanks to supervisors, colleagues, and family for support during the thesis journey."]},
    {"title": "Questions", "content": [
        "Feel free to ask about any aspect of the thesis or implementation!"]}
]

# Function to add slides
def add_slide_with_content(presentation, title, content):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear default content

    for line in content:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        paragraph.alignment = PP_ALIGN.LEFT

# Create slides based on content
for slide_info in slides_content:
    add_slide_with_content(presentation, slide_info["title"], slide_info["content"])

# Save the presentation
presentation.save("Bachelor_Thesis_Presentation.pptx")
print("Presentation saved as 'Bachelor_Thesis_Presentation.pptx'")
