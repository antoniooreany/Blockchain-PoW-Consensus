#   Copyright (c) 2024, Anton Gorshkov
#   All rights reserved.
#   This code is for a main.py and its unit tests.
#   For any questions or concerns, please contact Anton Gorshkov at antoniooreany@gmail.com

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF
import os

# Create a directory for extracted images
output_dir = "./data/extracted_images"
os.makedirs(output_dir, exist_ok=True)

# Function to extract images from a PDF
def extract_images_from_pdf(pdf_path, output_dir):
    doc = fitz.open(pdf_path)
    image_files = []
    for i in range(len(doc)):
        for img_index, img in enumerate(doc[i].get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_path = os.path.join(output_dir, f"page_{i + 1}_img_{img_index + 1}.{image_ext}")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            image_files.append(image_path)
    return image_files

# Path to the uploaded PDF
# pdf_path = "/mnt/data/Bachelor_Thesis_CTS_Gorshkov_Anton_3130111_signed.pdf"
pdf_path = "./Bachelor_Thesis_CTS_Gorshkov_Anton_3130111_signed.pdf"

# Extract images
extracted_images = extract_images_from_pdf(pdf_path, output_dir)

# Slide content
slides_content = [
    {"title": "Blockchain Technology", "content": [
        "Bachelor Thesis: Consensus Mechanism Proof-of-Work",
        "Author: Anton Gorshkov",
        "Supervisors: Prof. Dr. rer. nat. Martin Severin, Prof. Dr.-Ing. Philipp Graf",
        "Submission Date: 30.11.2024"], "image": None},
    {"title": "Introduction", "content": [
        "Motivation: Blockchain as a disruptive innovation.",
        "Problem Statement: Challenges in energy, scalability, and centralization.",
        "Goal: Develop a blockchain prototype addressing these challenges."], "image": extracted_images[0] if extracted_images else None},
    {"title": "Research Questions", "content": [
        "How does floating-point bit difficulty improve performance and stability?",
        "What are the trade-offs between simplicity and scalability?",
        "How does dynamic difficulty adjustment influence block mining time?"], "image": None},
    {"title": "Implementation", "content": [
        "Core Components: Block and Blockchain classes.",
        "Features: Dynamic difficulty adjustment, Proof-of-Work integration."], "image": extracted_images[1] if len(extracted_images) > 1 else None},
]

# Create a PowerPoint presentation
presentation = Presentation()

# Function to add slide with content and optional image
def add_slide_with_content(presentation, title, content, image_path=None):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = title
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear default content

    for line in content:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        paragraph.alignment = PP_ALIGN.LEFT

    # Add image if available
    if image_path and os.path.exists(image_path):
        left = Inches(5)  # Position image to the right
        top = Inches(1)
        slide.shapes.add_picture(image_path, left, top, width=Inches(4))

# Generate slides
for slide_info in slides_content:
    add_slide_with_content(
        presentation, slide_info["title"], slide_info["content"], slide_info.get("image")
    )

# Save the presentation
output_path = "./Enhanced_Bachelor_Thesis_Presentation_With_Images.pptx"
presentation.save(output_path)

output_path
